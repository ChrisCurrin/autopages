from pathlib import Path
import pptx


def analyze_ppt(template_name):
    """Take the input file and analyze the structure.
    The output file contains marked up information to make it easier
    for generating future powerpoint templates.

    :param template_name: Full path of the template file.

    """
    prs = pptx.Presentation(template_name)
    # Each powerpoint file has multiple layouts
    # Loop through them all and  see where the various elements are
    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = "Title for Layout {}".format(index)
        except AttributeError:
            print("No Title for Layout {}".format(index))
        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is just a special placeholder
                try:
                    if "Title" not in shape.text:
                        shape.text = "Placeholder index:{} type:{}".format(
                            phf.idx, shape.name
                        )
                except AttributeError:
                    print("{} has no text attribute".format(phf.type))
                print("{} {}".format(phf.idx, shape.name))
    # add markup to template name, taking into account the filetype
    prs.save(str(template_name).replace(".pptx", "-markup.pptx"))


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="Analyze a powerpoint template"
    )
    parser.add_argument(
        "template",
        type=argparse.FileType("r"),
        help="Powerpoint file used as the template",
    )
    args = parser.parse_args()
    template_name = Path(args.template.name).resolve()
    analyze_ppt(template_name)
