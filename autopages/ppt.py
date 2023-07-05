from __future__ import annotations
import logging
from datetime import datetime
from pathlib import Path
from zipfile import BadZipFile


import pptx


logger = logging.getLogger(__name__)


default_template = Path(__file__).parent / ".." / "assets" / "template.pptx"

FIGSIZE = (7.13, 6.50)  # according to template
FIGSIZE_WIDE = (7.13, 1.8)


def create_ppt(
    template_name: str | Path,
    output_name: str | Path,
    data: dict[str, dict[str, str]] | dict[str, str],
    no_titles: bool = False,
):
    """Take the input powerpoint file and use it as the template for the output
    file.
    """
    logger.info(f"Creating powerpoint report '{output_name}' ...")
    try:
        prs = pptx.Presentation(template_name)
    except BadZipFile as err:
        logger.error(
            "Template file is not a valid powerpoint file. "
            "Try `git lfs pull` if the pptx file is plain text"
        )
        raise err

    for index, (page_title, page_data) in enumerate(data.items()):
        if not isinstance(page_data, dict) and isinstance(page_data, str):
            page_data = {
                "content": page_data,
                "title": page_title,
                "layout": 0,
            }

        if index >= len(prs.slides):
            if "layout" not in page_data:
                page_data["layout"] = 0
            slide = prs.slides.add_slide(
                prs.slide_layouts[int(str(page_data["layout"]))]
            )
        else:
            # use the existing slide
            slide = prs.slides[index]

        title = slide.shapes.title
        placeholders = slide.placeholders

        ###########################################################################
        # Get placeholders
        ###########################################################################

        pics = [
            p
            for p in placeholders
            if isinstance(
                p, pptx.shapes.placeholder.PicturePlaceholder  # type: ignore
            )
        ]
        content_block = [
            p
            for p in placeholders
            if isinstance(
                p, pptx.shapes.placeholder.SlidePlaceholder  # type: ignore
            )
            and p not in pics
            and p != title
        ]

        content = page_data["content"]

        if isinstance(content, str):
            content = [content]

        if len(content) > len(content_block):
            raise ValueError(
                "More content than content blocks. Please check the template"
            )

        date_placeholder = [
            p for p in placeholders if "date" in p.name.lower()
        ]

        slide_num_placeholder = [
            p for p in placeholders if "slide number" in p.name.lower()
        ]

        if len(date_placeholder) == 0 and len(slide_num_placeholder) == 0:
            # get from master layout
            # order is Date | Footer | Slide Number
            master_placeholders = [
                p for p in prs.slide_layouts.parent.placeholders
            ]
            for p in master_placeholders:
                slide.shapes.clone_placeholder(p)

            # re-get placeholders now that they have been placed
            date_placeholder = [
                p for p in placeholders if "date" in p.name.lower()
            ]

            slide_num_placeholder = [
                p for p in placeholders if "slide number" in p.name.lower()
            ]

        ###########################################################################
        # add content to the placeholders
        ###########################################################################
        page_data_title = page_data.get("title")
        if title is not None and not no_titles:
            for ignore in ["Unnamed", "None"]:
                if ignore in page_data_title:
                    break
            else:
                # no breaking
                title.text = page_data["title"]

        # TODO: add support for adding figures
        # # Add the figures to the slide
        # for pic, fig in zip(
        #     sorted(pics, key=lambda p: p.shape_id, reverse=True), content
        # ):
        #     pic.insert_picture(fig)

        # add content
        for block, text in zip(content_block, content):
            block.text = text

        # Add the date to the slide
        #   we do this manually as the timezone may be different when converting
        if len(date_placeholder) > 0:
            t = datetime.now().strftime("%d/%m/%Y %H:%M")
            for ph in date_placeholder:
                ph.text = f"generated on {t}"

        # Add the total number of pages to the slide number
        if len(slide_num_placeholder) > 0:
            for ph in slide_num_placeholder:
                ph.text = f"page {index+1}/{len(data)}"
    ###########################################################################
    # save the presentation (make sure directory exists)
    ###########################################################################

    if "content" in str(output_name):
        output_name = str(output_name).replace(
            "content", page_data["content"]
        )
    if "title" in str(output_name):
        output_name = str(output_name).replace(
            "title", page_data["title"]
        )

    Path(output_name).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_name)
